# =========================================================
# poc_bearing_predictive_maintenance.py
# Predictive Maintenance – Bearing RUL & Degradation
# =========================================================

import os
import json
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")  # чтобы работало в Docker / без GUI
import matplotlib.pyplot as plt

from scipy.stats import kurtosis, skew
from sklearn.ensemble import IsolationForest
from sklearn.preprocessing import MinMaxScaler
from sklearn.linear_model import LinearRegression

from pptx import Presentation
from pptx.util import Inches

# =========================
# CONFIG
# =========================
# Если хочешь – можно сделать условие под Windows/Docker, как раньше
DATA_DIR = "data/1st_test"
OUTPUT_DIR = "outputs"

os.makedirs(OUTPUT_DIR, exist_ok=True)

TRAIN_RATIO = 0.3             # сколько считаем "здоровыми" в начале
MINUTES_PER_FILE = 10         # шаг по времени между файлами
WINDOW_SIZE = 20              # сглаживание признаков по окну

# уровни health для зон (в процентах)
HEALTH_GREEN = 70
HEALTH_YELLOW = 40
HEALTH_ORANGE = 20

# =========================
# LOAD + FEATURE EXTRACTION
# =========================
features = []
files = sorted(os.listdir(DATA_DIR))

for fname in files:
    if fname.startswith("."):
        continue

    path = os.path.join(DATA_DIR, fname)

    try:
        signal = np.loadtxt(path)
    except Exception as e:
        print(f"⚠ Не удалось загрузить {fname}: {e}")
        continue

    signal = np.ravel(signal)
    if signal.size == 0:
        print(f"⚠ Пустой сигнал в {fname}, пропуск")
        continue

    rms = np.sqrt(np.mean(signal ** 2))
    std = np.std(signal)
    k = kurtosis(signal, fisher=True, bias=False)
    s = skew(signal, bias=False)
    if np.isnan(k):
        k = 0.0
    if np.isnan(s):
        s = 0.0
    p2p = np.ptp(signal)

    features.append([rms, std, k, s, p2p])

if not features:
    raise RuntimeError("Нет валидных файлов/сигналов. Проверь DATA_DIR.")

df = pd.DataFrame(features, columns=["rms", "std", "kurtosis", "skew", "p2p"])

# =========================
# SLIDING WINDOW (SMOOTHING)
# =========================
df_sw = df.rolling(WINDOW_SIZE, min_periods=WINDOW_SIZE).mean()
df_sw = df_sw.bfill()  # заполняем первые значения

# =========================
# TRAIN ANOMALY MODEL
# =========================
train_len = int(len(df_sw) * TRAIN_RATIO)
if train_len < 1:
    raise RuntimeError("Слишком мало данных для TRAIN_RATIO. Увеличь количество файлов.")

X_train = df_sw.iloc[:train_len]
X_all = df_sw

model = IsolationForest(
    n_estimators=300,
    contamination=0.05,
    random_state=42
)
model.fit(X_train)

# anomaly_score: больше = более аномально (хуже)
anomaly_score = -model.decision_function(X_all)

# =========================
# HEALTH INDEX (0–100, 100 = хорошо)
# =========================
# Нормируем anomaly_score в [0, 100], где 100 = максимум аномалии
scaler = MinMaxScaler((0, 100))
score_norm = scaler.fit_transform(anomaly_score.reshape(-1, 1)).flatten()

# Health: 100 = отлично, 0 = ужасно
health = 100 - score_norm

# =========================
# DEGRADATION EVENTS (ПО HEALTH)
# =========================
def first_crossing_below(arr, threshold):
    idx = np.where(arr < threshold)[0]
    return int(idx[0]) if len(idx) > 0 else None

first_warning_idx = first_crossing_below(health, HEALTH_GREEN)   # начало деградации
critical_idx = first_crossing_below(health, HEALTH_ORANGE)       # критическая зона
failure_idx_observed = len(health) - 1                           # последний известный файл

# =========================
# RUL: ПРОГНОЗ ДО ОТКАЗА
# =========================
time_idx = np.arange(len(health)).reshape(-1, 1)

reg = LinearRegression()
reg.fit(time_idx, health)

trend = reg.predict(time_idx)
residuals = health - trend
sigma = np.std(residuals)

coef = reg.coef_[0]
intercept = reg.intercept_

# Прогнозируем момент, когда trend ≈ 0 (здоровье падает до 0)
if abs(coef) < 1e-6:
    # почти нет тренда – считаем, что предсказать отказ нельзя, берем последний индекс
    failure_pred_idx = failure_idx_observed
else:
    failure_pred_idx = int(-intercept / coef)
    # не позволяем "прогнозу" быть раньше последнего наблюдения
    failure_pred_idx = max(failure_pred_idx, failure_idx_observed)

# Переводим шаги в дни: сколько осталось от каждого момента до predicted failure
rul_steps = failure_pred_idx - time_idx.flatten()
rul_days = rul_steps * MINUTES_PER_FILE / (60 * 24)

# Грубая 95% CI по индексу (для демонстрации, а не строгой статистики)
ci_lower_idx = failure_pred_idx - 1.96 * sigma
ci_upper_idx = failure_pred_idx + 1.96 * sigma
ci_lower_days = ci_lower_idx * MINUTES_PER_FILE / (60 * 24)
ci_upper_days = ci_upper_idx * MINUTES_PER_FILE / (60 * 24)

# =========================
# METRICS – ПРЕДИКТИВНОЕ ОБСЛУЖИВАНИЕ
# =========================
metrics = {}

# Текущее состояние (в начале данных, индекс 0)
metrics["RUL now (days from start)"] = round(rul_days[0], 2)

# Когда начались первые признаки деградации
if first_warning_idx is not None:
    metrics["First warning index"] = int(first_warning_idx)
    metrics["First warning health"] = round(health[first_warning_idx], 2)
    metrics["RUL at first warning (days)"] = round(rul_days[first_warning_idx], 2)
else:
    metrics["First warning index"] = None
    metrics["First warning health"] = None
    metrics["RUL at first warning (days)"] = None

# Когда система стала критической
if critical_idx is not None:
    metrics["Critical index"] = int(critical_idx)
    metrics["Critical health"] = round(health[critical_idx], 2)
    metrics["RUL at critical (days)"] = round(rul_days[critical_idx], 2)
else:
    metrics["Critical index"] = None
    metrics["Critical health"] = None
    metrics["RUL at critical (days)"] = None

# Наблюдаемый и предсказанный отказ
metrics["Last observed index"] = int(failure_idx_observed)
metrics["Predicted failure index"] = int(failure_pred_idx)
metrics["Estimated failure RUL from start (days)"] = round(rul_days[0], 2)
metrics["RUL CI lower (days)"] = round(ci_lower_days, 2)
metrics["RUL CI upper (days)"] = round(ci_upper_days, 2)

with open(f"{OUTPUT_DIR}/metrics_predective.json", "w") as f:
    json.dump(metrics, f, indent=2, ensure_ascii=False)

# =========================
# VISUALIZATION
# =========================
plt.figure(figsize=(14, 6))

plt.plot(health, label="Health (100=good, 0=bad)", color="black")
plt.plot(trend, "--", label="Degradation trend", color="blue")

plt.axhspan(HEALTH_GREEN, 100, color="green", alpha=0.15, label="Healthy zone")
plt.axhspan(HEALTH_YELLOW, HEALTH_GREEN, color="yellow", alpha=0.2, label="Warning zone")
plt.axhspan(HEALTH_ORANGE, HEALTH_YELLOW, color="orange", alpha=0.25, label="Degradation zone")
plt.axhspan(0, HEALTH_ORANGE, color="red", alpha=0.3, label="Critical zone")

if first_warning_idx is not None:
    plt.axvline(first_warning_idx, linestyle="--", color="orange", label="First warning")
if critical_idx is not None:
    plt.axvline(critical_idx, linestyle="--", color="red", label="Critical")
plt.axvline(failure_pred_idx, linestyle="--", color="purple", label="Predicted failure")

plt.title("Bearing Degradation & RUL – Predictive Maintenance")
plt.xlabel("Time index")
plt.ylabel("Health (%)")
plt.legend()
plt.grid(True)

plot_path = f"{OUTPUT_DIR}/health_rul_predective.png"
plt.savefig(plot_path, dpi=300, bbox_inches="tight")
plt.close()

# =========================
# POWERPOINT REPORT
# =========================
prs = Presentation()

# Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Bearing Predictive Maintenance"
slide.placeholders[1].text = (
    "PoC – Degradation detection & RUL estimation\n\n"
    f"RUL now: ~{metrics['RUL now (days from start)']} days\n"
    f"RUL at first warning: ~{metrics['RUL at first warning (days)']} days\n"
    f"RUL at critical: ~{metrics['RUL at critical (days)']} days\n"
    f"Estimated failure RUL (from start): ~{metrics['Estimated failure RUL from start (days)']} days"
)

# Architecture
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Edge → Cloud Architecture (PoC vision)"
slide.placeholders[1].text = (
    "Edge:\n"
    "- Vibration sensor acquisition\n"
    "- Feature extraction (RMS, STD, kurtosis, skew, p2p)\n"
    "- Sliding window smoothing\n\n"
    "Cloud:\n"
    "- Anomaly detection (Isolation Forest)\n"
    "- Health Index (0–100)\n"
    "- Degradation trend & RUL prediction\n"
    "- Visualization & reporting"
)

# Result
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Degradation & RUL Result"
slide.shapes.add_picture(
    plot_path,
    Inches(1),
    Inches(1.5),
    width=Inches(8)
)

# Metrics slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Key Predictive Metrics"
slide.placeholders[1].text = "\n".join(
    [f"{k}: {v}" for k, v in metrics.items()]
)

prs.save(f"{OUTPUT_DIR}/Bearing_Predictive_Report.pptx")

print("\n✅ Predictive maintenance PoC successfully completed")
print("Artifacts:")
print(f"- {plot_path}")
print(f"- {OUTPUT_DIR}/metrics_predective.json")
print(f"- {OUTPUT_DIR}/Bearing_Predictive_Report.pptx")
